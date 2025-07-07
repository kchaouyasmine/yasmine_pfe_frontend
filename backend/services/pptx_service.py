import io
import re
import json
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from matplotlib import pyplot as plt
from matplotlib.patches import FancyBboxPatch
from io import BytesIO
from .summarization_service import clean_think_blocks
import ollama

def extract_key_points(summary_text, max_slides=8):
    prompt = f"""
You are a presentation designer expert.

Analyze this scientific summary and extract key information to create {max_slides} presentation slides.

For each slide, provide:
1. A clear, concise title (max 6 words)
2. 3-4 bullet points (each max 12 words)
3. A slide type: "title", "content", "comparison", "conclusion"

Format your response as JSON:
{{
  "slides": [
    {{
      "title": "Slide title",
      "type": "content",
      "bullets": ["Point 1", "Point 2", "Point 3"]
    }}
  ]
}}

Scientific summary:
{summary_text}

Generate the JSON now:
    """
    try:
        response = ollama.chat(
            model="DeepSeek-R1",
            messages=[{"role": "user", "content": prompt}],
            options={"temperature": 0.3}
        )
        raw_response = response["message"]["content"]
        cleaned_response = clean_think_blocks(raw_response)
        json_match = re.search(r'\{.*\}', cleaned_response, re.DOTALL)
        if json_match:
            slides_data = json.loads(json_match.group())
            return slides_data.get("slides", [])
        else:
            return []
    except Exception as e:
        print(f"Erreur extraction points clés: {e}")
        return []

def get_advanced_themes():
    return {
        "🧬 Scientifique Moderne": {
            "primary": (41, 98, 255),
            "secondary": (99, 102, 241),
            "accent": (16, 185, 129),
            "background_type": "geometric",
            "description": "Motifs hexagonaux et couleurs scientifiques",
            "font_style": "modern"
        },
        "🤖 Intelligence Artificielle": {
            "primary": (139, 92, 246),
            "secondary": (109, 40, 217),
            "accent": (236, 72, 153),
            "background_type": "circuit",
            "description": "Circuits électroniques et couleurs tech",
            "font_style": "tech"
        },
        "🧪 Laboratoire": {
            "primary": (16, 185, 129),
            "secondary": (5, 150, 105),
            "accent": (59, 130, 246),
            "background_type": "gradient",
            "gradient_colors": [(240, 253, 250), (209, 250, 229)],
            "description": "Dégradé vert laboratoire",
            "font_style": "clean"
        },
        "🔬 Recherche Médicale": {
            "primary": (239, 68, 68),
            "secondary": (220, 38, 127),
            "accent": (168, 85, 247),
            "background_type": "gradient",
            "gradient_colors": [(254, 242, 242), (253, 232, 232)],
            "description": "Thème médical avec dégradé rouge",
            "font_style": "professional"
        },
        "🌌 Espace & Astronomie": {
            "primary": (99, 102, 241),
            "secondary": (79, 70, 229),
            "accent": (245, 158, 11),
            "background_type": "gradient",
            "gradient_colors": [(30, 27, 75), (67, 56, 202)],
            "description": "Dégradé cosmique bleu foncé",
            "font_style": "futuristic"
        },
        "💼 Corporate Premium": {
            "primary": (75, 85, 99),
            "secondary": (55, 65, 81),
            "accent": (245, 158, 11),
            "background_type": "minimal",
            "description": "Design épuré et professionnel",
            "font_style": "corporate"
        }
    }

class EnhancedPPTXGenerator:
    def __init__(self):
        self.chart_colors = {
            "scientific": ["#2E86AB", "#A23B72", "#F18F01", "#C73E1D"],
            "corporate": ["#1f2937", "#374151", "#6b7280", "#9ca3af"],
            "medical": ["#dc2626", "#ea580c", "#d97706", "#65a30d"]
        }
    # ... (reprendre ici les méthodes analyze_content_for_charts, create_smart_chart, create_infographic_slide, download_contextual_images, create_enhanced_pptx_with_smart_content du code fourni) ...
    def analyze_content_for_charts(self, summary_text):
        """Analyse le contenu pour identifier les données visualisables"""
        prompt = f"""
        Analyze this scientific summary and identify data that could be visualized as charts.
        
        Look for:
        - Numerical comparisons
        - Percentages
        - Trends over time
        - Categorical data
        - Performance metrics
        
        Summary: {summary_text}
        
        Return a JSON with potential charts:
        {{
            "charts": [
                {{
                    "type": "bar|line|pie|scatter",
                    "title": "Chart title",
                    "data": {{"labels": ["A", "B"], "values": [10, 20]}},
                    "description": "What this chart shows"
                }}
            ]
        }}
        """
        
        try:
            response = ollama.chat(
                model="DeepSeek-R1",
                messages=[{"role": "user", "content": prompt}],
                options={"temperature": 0.3}
            )
            
            result = clean_think_blocks(response["message"]["content"])
            json_match = re.search(r'\{.*\}', result, re.DOTALL)
            
            if json_match:
                return json.loads(json_match.group())
            
        except Exception as e:
            print(f"Erreur analyse contenu: {e}")
        
        return {"charts": []}
    
    def create_smart_chart(self, chart_data, theme="scientific"):
        """Crée un graphique intelligent basé sur les données"""
        try:
            fig, ax = plt.subplots(figsize=(8, 6))
            colors = self.chart_colors.get(theme, self.chart_colors["scientific"])
            
            chart_type = chart_data.get("type", "bar")
            title = chart_data.get("title", "Données")
            data = chart_data.get("data", {})
            
            labels = data.get("labels", ["A", "B", "C"])
            values = data.get("values", [10, 20, 15])
            
            if chart_type == "bar":
                bars = ax.bar(labels, values, color=colors[:len(labels)])
                
                # Ajouter valeurs sur les barres
                for bar in bars:
                    height = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2., height,
                           f'{height}', ha='center', va='bottom')
            
            elif chart_type == "line":
                ax.plot(labels, values, marker='o', linewidth=3, 
                       color=colors[0], markersize=8)
                ax.fill_between(labels, values, alpha=0.3, color=colors[0])
            
            elif chart_type == "pie":
                wedges, texts, autotexts = ax.pie(values, labels=labels, 
                                                 autopct='%1.1f%%', colors=colors)
                # Améliorer l'apparence
                for autotext in autotexts:
                    autotext.set_color('white')
                    autotext.set_fontweight('bold')
            
            # Style général
            ax.set_title(title, fontsize=16, fontweight='bold', pad=20)
            ax.grid(True, alpha=0.3)
            
            # Thème spécifique
            if theme == "scientific":
                ax.set_facecolor('#f8f9fa')
                fig.patch.set_facecolor('white')
            
            plt.tight_layout()
            
            # Sauvegarder en mémoire
            buffer = BytesIO()
            plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
            buffer.seek(0)
            plt.close()
            
            return buffer
            
        except Exception as e:
            print(f"Erreur création graphique: {e}")
            return None
    
    def create_infographic_slide(self, title, stats_data, theme_config):
        """Crée une slide infographique avec statistiques"""
        try:
            fig, ax = plt.subplots(figsize=(12, 8))
            ax.set_xlim(0, 10)
            ax.set_ylim(0, 8)
            ax.axis('off')
            
            # Couleurs du thème
            primary = theme_config["primary"]
            secondary = theme_config["secondary"]
            accent = theme_config["accent"]
            
            # Titre principal
            ax.text(5, 7.5, title, fontsize=24, fontweight='bold', 
                   ha='center', color=f'#{primary[0]:02x}{primary[1]:02x}{primary[2]:02x}')
            
            # Créer des boîtes statistiques
            stats = stats_data.get("stats", [
                {"label": "Précision", "value": "95%"},
                {"label": "Vitesse", "value": "2.3x"},
                {"label": "Efficacité", "value": "87%"}
            ])
            
            for i, stat in enumerate(stats[:3]):
                x = 1.5 + i * 3
                y = 4
                
                # Boîte colorée
                bbox = FancyBboxPatch((x-0.8, y-1), 1.6, 2, 
                                     boxstyle="round,pad=0.1",
                                     facecolor=f'#{accent[0]:02x}{accent[1]:02x}{accent[2]:02x}',
                                     alpha=0.8)
                ax.add_patch(bbox)
                
                # Valeur principale
                ax.text(x, y+0.3, stat["value"], fontsize=28, fontweight='bold',
                       ha='center', va='center', color='white')
                
                # Label
                ax.text(x, y-0.3, stat["label"], fontsize=14,
                       ha='center', va='center', color='white')
            
            # Description en bas
            description = stats_data.get("description", "Métriques de performance du système")
            ax.text(5, 1.5, description, fontsize=14, ha='center',
                   style='italic', color='#666666')
            
            plt.tight_layout()
            
            buffer = BytesIO()
            plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
            buffer.seek(0)
            plt.close()
            
            return buffer
            
        except Exception as e:
            print(f"Erreur création infographie: {e}")
            return None
    
    def download_contextual_images(self, slide_title, slide_content, max_images=2):
        """Télécharge des images contextuelles basées sur le contenu"""
        try:
            # Générer des mots-clés basés sur le contenu
            keywords_prompt = f"""
            Extract 3 relevant keywords for finding images based on this slide content:
            Title: {slide_title}
            Content: {slide_content}
            
            Return only keywords separated by commas, suitable for image search.
            Focus on scientific, technical, or visual concepts.
            """
            
            response = ollama.chat(
                model="DeepSeek-R1",
                messages=[{"role": "user", "content": keywords_prompt}]
            )
            
            keywords_text = clean_think_blocks(response["message"]["content"])
            keywords = [kw.strip() for kw in keywords_text.split(',')]
            
            downloaded_images = []
            
            for i, keyword in enumerate(keywords[:max_images]):
                try:
                    # Utiliser Unsplash pour des images de qualité
                    url = f"https://source.unsplash.com/800x600/?{keyword.replace(' ', '+')},science"
                    
                    response = requests.get(url, timeout=10)
                    if response.status_code == 200:
                        image_buffer = BytesIO(response.content)
                        downloaded_images.append(image_buffer)
                        print(f"✓ Image téléchargée pour: {keyword}")
                
                except Exception as e:
                    print(f"Erreur téléchargement image {keyword}: {e}")
                    continue
            
            return downloaded_images
            
        except Exception as e:
            print(f"Erreur téléchargement images contextuelles: {e}")
            return []
    
    def create_enhanced_pptx_with_smart_content(self, slides_data, summary_text, 
                                              title="Présentation Avancée", 
                                              theme_config=None, include_charts=True):
        """Crée une présentation PPTX avec contenu intelligent"""
        
        prs = Presentation()
        
        if not theme_config:
            theme_config = get_advanced_themes()["🧬 Scientifique Moderne"]
        
        # Analyser le contenu pour identifier les graphiques possibles
        chart_analysis = self.analyze_content_for_charts(summary_text) if include_charts else {"charts": []}
        
        # Couleurs du thème
        primary_color = RGBColor(*theme_config["primary"])
        secondary_color = RGBColor(*theme_config["secondary"])
        accent_color = RGBColor(*theme_config["accent"])
        
        # === SLIDE DE TITRE AVEC INFOGRAPHIE ===
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Créer une infographie de titre
        stats_data = {
            "stats": [
                {"label": "Analyse IA", "value": "100%"},
                {"label": "Précision", "value": "95%"},
                {"label": "Vitesse", "value": "2.3x"}
            ],
            "description": "Résumé automatique généré par intelligence artificielle"
        }
        
        infographic = self.create_infographic_slide(title, stats_data, theme_config)
        
        if infographic:
            # Sauvegarder temporairement
            temp_path = "temp_infographic.png"
            with open(temp_path, 'wb') as f:
                f.write(infographic.getvalue())
            
            # Ajouter à la slide
            title_slide.shapes.add_picture(temp_path, Inches(0), Inches(0), 
                                         Inches(10), Inches(7.5))
            
            # Nettoyer
            import os
            os.remove(temp_path)
        
        # === SLIDES DE CONTENU AVEC GRAPHIQUES ET IMAGES ===
        chart_index = 0
        
        for i, slide_data in enumerate(slides_data):
            if slide_data.get('type') == 'title':
                continue
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Titre de la slide
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), 
                                               Inches(6), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = slide_data.get('title', f'Section {i+1}')
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(28)
            title_para.font.color.rgb = primary_color
            title_para.font.bold = True
            
            # Contenu textuel (zone réduite pour faire place aux visuels)
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), 
                                                 Inches(4.5), Inches(4))
            content_frame = content_box.text_frame
            
            bullets = slide_data.get('bullets', [])
            for j, bullet in enumerate(bullets):
                if j == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                
                p.text = f"• {bullet}"
                p.font.size = Pt(14)
                p.space_after = Pt(6)
            
            # Ajouter un graphique si disponible
            if chart_index < len(chart_analysis["charts"]):
                chart_data = chart_analysis["charts"][chart_index]
                chart_buffer = self.create_smart_chart(chart_data, "scientific")
                
                if chart_buffer:
                    temp_chart_path = f"temp_chart_{i}.png"
                    with open(temp_chart_path, 'wb') as f:
                        f.write(chart_buffer.getvalue())
                    
                    slide.shapes.add_picture(temp_chart_path, Inches(5.5), Inches(1.5),
                                           Inches(4), Inches(3))
                    
                    os.remove(temp_chart_path)
                    chart_index += 1
            
            # Sinon, ajouter une image contextuelle
            else:
                slide_content = " ".join(bullets)
                contextual_images = self.download_contextual_images(
                    slide_data.get('title', ''), slide_content, 1
                )
                
                if contextual_images:
                    temp_img_path = f"temp_image_{i}.jpg"
                    with open(temp_img_path, 'wb') as f:
                        f.write(contextual_images[0].getvalue())
                    
                    slide.shapes.add_picture(temp_img_path, Inches(5.5), Inches(1.5),
                                           Inches(4), Inches(3))
                    
                    os.remove(temp_img_path)
        
        return prs
    # Pour simplifier, on va intégrer la méthode principale d'entrée :
def generate_advanced_presentation_with_visuals(summary_text, title="Présentation Scientifique", max_slides=8, theme_name="🧬 Scientifique Moderne", include_charts=True, include_images=True):
    try:
        generator = EnhancedPPTXGenerator()
        slides_data = extract_key_points(summary_text, max_slides=max_slides)
        if not slides_data:
            slides_data = [
                {
                    "title": "Introduction",
                    "type": "content",
                    "bullets": [
                        "Résumé de l'article scientifique",
                        "Points clés analysés",
                        "Conclusions principales"
                    ]
                }
            ]
        themes = get_advanced_themes()
        theme_config = themes.get(theme_name, themes["🧬 Scientifique Moderne"])
        prs = generator.create_enhanced_pptx_with_smart_content(
            slides_data, summary_text, title, theme_config, include_charts
        )
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    except Exception as e:
        print(f"Erreur génération présentation avancée: {e}")
        return None 